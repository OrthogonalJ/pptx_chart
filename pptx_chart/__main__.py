import collections
import argparse
import copy

import pptx
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
import pandas as pd
from tqdm import tqdm
import numpy as np

from pptx_chart.not_found_error import NotFoundError

LEGEND_POSITIONS = {
    'bottom': pptx.enum.chart.XL_LEGEND_POSITION.BOTTOM,
    'corner': pptx.enum.chart.XL_LEGEND_POSITION.CORNER,
    'left': pptx.enum.chart.XL_LEGEND_POSITION.LEFT,
    'right': pptx.enum.chart.XL_LEGEND_POSITION.RIGHT,
    'top': pptx.enum.chart.XL_LEGEND_POSITION.TOP
}

TICK_LABEL_POSITION = {
    'high': pptx.enum.chart.XL_TICK_LABEL_POSITION.HIGH,
    'low': pptx.enum.chart.XL_TICK_LABEL_POSITION.LOW,
    'next_to_axis': pptx.enum.chart.XL_TICK_LABEL_POSITION.NEXT_TO_AXIS,
    'none': pptx.enum.chart.XL_TICK_LABEL_POSITION.NONE
}

LINE_DASH_STYLE = {
    'dash': pptx.enum.dml.MSO_LINE.DASH,
    'dash_dot': pptx.enum.dml.MSO_LINE.DASH_DOT,
    'dash_dot_dot': pptx.enum.dml.MSO_LINE.DASH_DOT_DOT,
    'long_dash': pptx.enum.dml.MSO_LINE.LONG_DASH,
    'long_dash_dot': pptx.enum.dml.MSO_LINE.LONG_DASH_DOT,
    'round_dot': pptx.enum.dml.MSO_LINE.ROUND_DOT,
    'solid': pptx.enum.dml.MSO_LINE.SOLID,
    'square_dot': pptx.enum.dml.MSO_LINE.SQUARE_DOT,
    'dash_style_mixed': pptx.enum.dml.MSO_LINE.DASH_STYLE_MIXED,
}

CHART_TYPE = {
    'three_d_area': XL_CHART_TYPE.THREE_D_AREA,
    'three_d_area_stacked': XL_CHART_TYPE.THREE_D_AREA_STACKED,
    'three_d_area_stacked_100': XL_CHART_TYPE.THREE_D_AREA_STACKED_100,
    'three_d_bar_clustered': XL_CHART_TYPE.THREE_D_BAR_CLUSTERED,
    'three_d_bar_stacked': XL_CHART_TYPE.THREE_D_BAR_STACKED,
    'three_d_bar_stacked_100': XL_CHART_TYPE.THREE_D_BAR_STACKED_100,
    'three_d_column': XL_CHART_TYPE.THREE_D_COLUMN,
    'three_d_column_clustered': XL_CHART_TYPE.THREE_D_COLUMN_CLUSTERED,
    'three_d_column_stacked': XL_CHART_TYPE.THREE_D_COLUMN_STACKED,
    'three_d_column_stacked_100': XL_CHART_TYPE.THREE_D_COLUMN_STACKED_100,
    'three_d_line': XL_CHART_TYPE.THREE_D_LINE,
    'three_d_pie': XL_CHART_TYPE.THREE_D_PIE,
    'three_d_pie_exploded': XL_CHART_TYPE.THREE_D_PIE_EXPLODED,
    'area': XL_CHART_TYPE.AREA,
    'area_stacked': XL_CHART_TYPE.AREA_STACKED,
    'area_stacked_100': XL_CHART_TYPE.AREA_STACKED_100,
    'bar_clustered': XL_CHART_TYPE.BAR_CLUSTERED,
    'bar_of_pie': XL_CHART_TYPE.BAR_OF_PIE,
    'bar_stacked': XL_CHART_TYPE.BAR_STACKED,
    'bar_stacked_100': XL_CHART_TYPE.BAR_STACKED_100,
    'bubble': XL_CHART_TYPE.BUBBLE,
    'bubble_three_d_effect': XL_CHART_TYPE.BUBBLE_THREE_D_EFFECT,
    'column_clustered': XL_CHART_TYPE.COLUMN_CLUSTERED,
    'column_stacked': XL_CHART_TYPE.COLUMN_STACKED,
    'column_stacked_100': XL_CHART_TYPE.COLUMN_STACKED_100,
    'cone_bar_clustered': XL_CHART_TYPE.CONE_BAR_CLUSTERED,
    'cone_bar_stacked': XL_CHART_TYPE.CONE_BAR_STACKED,
    'cone_bar_stacked_100': XL_CHART_TYPE.CONE_BAR_STACKED_100,
    'cone_col': XL_CHART_TYPE.CONE_COL,
    'cone_col_clustered': XL_CHART_TYPE.CONE_COL_CLUSTERED,
    'cone_col_stacked': XL_CHART_TYPE.CONE_COL_STACKED,
    'cone_col_stacked_100': XL_CHART_TYPE.CONE_COL_STACKED_100,
    'cylinder_bar_clustered': XL_CHART_TYPE.CYLINDER_BAR_CLUSTERED,
    'cylinder_bar_stacked': XL_CHART_TYPE.CYLINDER_BAR_STACKED,
    'cylinder_bar_stacked_100': XL_CHART_TYPE.CYLINDER_BAR_STACKED_100,
    'cylinder_col': XL_CHART_TYPE.CYLINDER_COL,
    'cylinder_col_clustered': XL_CHART_TYPE.CYLINDER_COL_CLUSTERED,
    'cylinder_col_stacked': XL_CHART_TYPE.CYLINDER_COL_STACKED,
    'cylinder_col_stacked_100': XL_CHART_TYPE.CYLINDER_COL_STACKED_100,
    'doughnut': XL_CHART_TYPE.DOUGHNUT,
    'doughnut_exploded': XL_CHART_TYPE.DOUGHNUT_EXPLODED,
    'line': XL_CHART_TYPE.LINE,
    'line_markers': XL_CHART_TYPE.LINE_MARKERS,
    'line_markers_stacked': XL_CHART_TYPE.LINE_MARKERS_STACKED,
    'line_markers_stacked_100': XL_CHART_TYPE.LINE_MARKERS_STACKED_100,
    'line_stacked': XL_CHART_TYPE.LINE_STACKED,
    'line_stacked_100': XL_CHART_TYPE.LINE_STACKED_100,
    'pie': XL_CHART_TYPE.PIE,
    'pie_exploded': XL_CHART_TYPE.PIE_EXPLODED,
    'pie_of_pie': XL_CHART_TYPE.PIE_OF_PIE,
    'pyramid_bar_clustered': XL_CHART_TYPE.PYRAMID_BAR_CLUSTERED,
    'pyramid_bar_stacked': XL_CHART_TYPE.PYRAMID_BAR_STACKED,
    'pyramid_bar_stacked_100': XL_CHART_TYPE.PYRAMID_BAR_STACKED_100,
    'pyramid_col': XL_CHART_TYPE.PYRAMID_COL,
    'pyramid_col_clustered': XL_CHART_TYPE.PYRAMID_COL_CLUSTERED,
    'pyramid_col_stacked': XL_CHART_TYPE.PYRAMID_COL_STACKED,
    'pyramid_col_stacked_100': XL_CHART_TYPE.PYRAMID_COL_STACKED_100,
    'radar': XL_CHART_TYPE.RADAR,
    'radar_filled': XL_CHART_TYPE.RADAR_FILLED,
    'radar_markers': XL_CHART_TYPE.RADAR_MARKERS,
    'stock_hlc': XL_CHART_TYPE.STOCK_HLC,
    'stock_ohlc': XL_CHART_TYPE.STOCK_OHLC,
    'stock_vhlc': XL_CHART_TYPE.STOCK_VHLC,
    'stock_vohlc': XL_CHART_TYPE.STOCK_VOHLC,
    'surface': XL_CHART_TYPE.SURFACE,
    'surface_top_view': XL_CHART_TYPE.SURFACE_TOP_VIEW,
    'surface_top_view_wireframe': XL_CHART_TYPE.SURFACE_TOP_VIEW_WIREFRAME,
    'surface_wireframe': XL_CHART_TYPE.SURFACE_WIREFRAME,
    'xy_scatter': XL_CHART_TYPE.XY_SCATTER,
    'xy_scatter_lines': XL_CHART_TYPE.XY_SCATTER_LINES,
    'xy_scatter_lines_no_markers': XL_CHART_TYPE.XY_SCATTER_LINES_NO_MARKERS,
    'xy_scatter_smooth': XL_CHART_TYPE.XY_SCATTER_SMOOTH,
    'xy_scatter_smooth_no_markers': XL_CHART_TYPE.XY_SCATTER_SMOOTH_NO_MARKERS
}


def handle_missing_chart_error(error, ignore_missing_charts):
    if ignore_missing_charts:
        print('WARNING:', error)
    else:
        raise error


def parse_bool(value):
    return str(value).lower() == 'true'


def parse_spec_cols(data, prefix, defaults):
    config = copy.deepcopy(defaults)
    for col_name in data.columns:
        if col_name.startswith(prefix):
            prop_name = col_name[len(prefix):]
            prop_value = data.loc[:, col_name].dropna().drop_duplicates().iloc[0]
            config[prop_name] = prop_value
    return config


def parse_y_specs(data):
    PREFIX = 'y.'
    specs = collections.defaultdict(dict)
    for col_name in data.columns:
        if col_name.startswith(PREFIX):
            series_key = col_name[len(PREFIX):].split('.')[0]
            prop_name = col_name[len(PREFIX) + len(series_key) + 1:]
            specs[series_key][prop_name] = data.loc[:, col_name].dropna().drop_duplicates().iloc[0]
    specs = list(specs.values())
    
    for spec in specs:
        spec['name'] = spec.get('name', spec['col'])
    
    return specs


def apply_axis_format(axis, spec):
    if 'title' in spec:
        axis.axis_title.text_frame.text = spec['title']
    if 'tick_font' in spec:
        axis.tick_labels.font.name = spec['tick_font']
    if 'tick_bold' in spec:
        axis.tick_labels.font.bold = parse_bool(spec['tick_bold'])
    if 'tick_italic' in spec:
        axis.tick_labels.font.italic = parse_bool(spec['tick_italic'])
    if 'tick_underline' in spec:
        axis.tick_labels.font.underline = parse_bool(spec['tick_underline'])
    if 'tick_color' in spec:
        axis.tick_labels.font.color.rgb = RGBColor.from_string(spec['tick_color'])
    if 'tick_color_brightness' in spec:
        axis.tick_labels.font.color.brightness = float(spec['tick_color_brightness'])
    if 'tick_size' in spec:
        axis.tick_labels.font.size = Pt(float(spec['tick_size']))
    if 'number_format' in spec:
        axis.tick_labels.number_format = spec['number_format']
    if 'tick_position' in spec:
        axis.tick_label_position = TICK_LABEL_POSITION[spec['tick_position']]


def clean_series_values(values):
    values = list(pd.to_numeric(values))
    values = [value if not (np.isnan(value) or np.isinf(value)) else None 
              for value in values]
    return values


def make_chart_data(data, x_spec, y_specs):
    chart_data = ChartData()

    categories = data.loc[:, x_spec['col']]
    if x_spec['type'] == 'date':
        categories = pd.to_datetime(categories)
    chart_data.categories = categories

    for spec in y_specs:
        col_name = spec['col']
        series_name = spec.get('name', col_name)
        values = clean_series_values(data.loc[:, col_name])
        chart_data.add_series(series_name, values)
    
    return chart_data


def parse_specs(data):
    y_specs = parse_y_specs(data)
    
    y_axis_spec = {}
    y_axis_spec = parse_spec_cols(data, 'y_axis.', y_axis_spec)

    x_spec = {
        'col': 'x',
        'type': 'string',
    }
    x_spec = parse_spec_cols(data, 'x_axis.', x_spec)

    legend_spec = {
        'enabled': 'true',
        'position': 'bottom'
    }
    legend_spec = parse_spec_cols(data, 'legend.', legend_spec)

    chart_spec = {
        'type': 'line',
        'width': '20.32',
        'height': '10.16'
    }
    chart_spec = parse_spec_cols(data, 'chart.', chart_spec)

    return {
        'y': y_specs,
        'x': x_spec,
        'legend': legend_spec,
        'chart': chart_spec,
        'y_axis': y_axis_spec
    }


def get_facet_iterator(data):
    facet_ids_col = data['facet.col'].dropna().drop_duplicates().iloc[0]
    facet_ids = data[facet_ids_col].dropna().drop_duplicates()
    for facet_id in tqdm(facet_ids):
        print('facet_id:', facet_id)
        facet_data = data.loc[data[facet_ids_col] == facet_id, :]
        yield facet_data


def format_chart(chart, specs):
    y_specs = specs['y']
    y_axis_spec = specs['y_axis']
    x_spec = specs['x']
    legend_spec = specs['legend']
    chart_spec = specs['chart']

    if 'title' in chart_spec:
        chart.chart_title.text_frame.text = chart_spec['title']
    if 'title_color' in chart_spec:
        chart.chart_title.text_frame.paragraphs[0].font.color.rgb = RGBColor.from_string(chart_spec['title_color'])
    if 'title_color_brightness' in chart_spec:
        chart.chart_title.text_frame.paragraphs[0].font.color.brightness = float(chart_spec['title_color_brightness'])
    if 'title_font' in chart_spec:
        chart.chart_title.text_frame.paragraphs[0].font.name = chart_spec['title_font']
    if 'title_size' in chart_spec:
        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(float(chart_spec['title_size']))

    category_axis = chart.category_axis
    apply_axis_format(category_axis, x_spec)

    value_axis = chart.value_axis
    apply_axis_format(value_axis, y_axis_spec)

    for i, spec in enumerate(y_specs):
        series = chart.series[i]
        series.smooth = parse_bool(spec.get('smooth', 'false'))
        if 'fill_color' in spec:
            series.format.fill.solid()
            series.format.fill.fore_color.rgb = RGBColor.from_string(spec['fill_color'])
        if 'fill_color_brightness' in spec:
            series.format.fill.fore_color.brightness = float(spec['fill_brightness'])
        if 'line_color' in spec:
            series.format.line.color.rgb = RGBColor.from_string(spec['line_color'])
        if 'line_color_brightness' in spec:
            series.format.line.color.brightness = float(spec['line_color_brightness'])
        if 'line_width' in spec:
            series.format.line.width = Pt(float(spec['line_width']))
        if 'line_dash' in spec:
            series.format.line.dash_style = LINE_DASH_STYLE[spec['line_dash']]

    legend_enabled = parse_bool(legend_spec['enabled'])
    chart.has_legend = legend_enabled
    if legend_enabled:
        chart.legend.position = LEGEND_POSITIONS[legend_spec['position']]


def make_chart(slide, data):
    specs = parse_specs(data)
    y_specs = specs['y']
    x_spec = specs['x']
    chart_spec = specs['chart']

    chart_data = make_chart_data(data, x_spec, y_specs)

    chart_shape = slide.shapes.add_chart(
        CHART_TYPE[chart_spec['type']],
        Cm(0),
        Cm(0),
        Cm(float(chart_spec['width'])),
        Cm(float(chart_spec['height'])),
        chart_data
    )
    if 'id' in chart_spec:
        chart_shape.name = chart_spec['id']
    chart = chart_shape.chart

    format_chart(chart, specs)


def make_facet_charts(slide, data):
    for facet_data in get_facet_iterator(data):
        make_chart(slide, facet_data)


def add_chart(output_file, data_file, slide_idx=None, input_file=None):
    if input_file is not None:
        presentation = pptx.Presentation(input_file)
    else:
        presentation =  pptx.Presentation()
        # NOTE: 1 = Title and content layout
        presentation.slides.add_slide(presentation.slide_layouts[1])
        slide_idx = 0
    slide = presentation.slides[slide_idx]

    data = pd.read_csv(data_file, dtype='str')

    if 'facet.col' in data.columns:
        make_facet_charts(slide, data)
    else:
        make_chart(slide, data)

    presentation.save(output_file)


def _update_chart(data, slide, shape_id, should_update_format):
    try:
        chart_shape = [shape for shape in slide.shapes if shape.name == shape_id][0]
    except IndexError:
        raise NotFoundError('Shape with id {} not found'.format(shape_id))
    
    specs = parse_specs(data)
    x_specs = specs['x']
    y_specs = specs['y']
    y_specs_indexed = {spec['name']: spec for spec in y_specs}
    # Ensure series order matches existing chart
    series_names = [series.name for series in chart_shape.chart.series]
    y_specs = [y_specs_indexed[name] for name in series_names]

    chart_data = make_chart_data(data, x_spec=x_specs, y_specs=y_specs)
    chart_shape.chart.replace_data(chart_data)

    if should_update_format:
        format_chart(chart_shape.chart, specs)


def update_facet_charts(data, slide, should_update_format, ignore_missing_charts):
    for facet_data in get_facet_iterator(data):
        specs = parse_specs(facet_data)
        shape_id = specs['chart']['id']
        try:
            _update_chart(facet_data, slide, shape_id, should_update_format)
        except NotFoundError as err:
            handle_missing_chart_error(err, ignore_missing_charts)


def update_chart(
        input_file, 
        data_file, 
        slide_idx, 
        shape_id=None, 
        output_file=None, 
        should_update_format=False,
        ignore_missing_charts=False):
    output_file = output_file if output_file is not None else input_file
    presentation = pptx.Presentation(input_file)
    slide = presentation.slides[slide_idx]

    data = pd.read_csv(data_file, dtype='str')

    if 'facet.col' in data.columns:
        update_facet_charts(data, slide, should_update_format, ignore_missing_charts)
    else:
        if shape_id is None:
            raise ValueError('Argument shape_id is required')
        try:
            _update_chart(data, slide, shape_id, should_update_format)
        except NotFoundError as err:
            handle_missing_chart_error(err, ignore_missing_charts)

    presentation.save(output_file)


def main():
    arg_parser = argparse.ArgumentParser(prog='pptx_chart')
    arg_parser.add_argument('-o', '--output_file', help='pptx file to write to. Defaults to -i/--input_file when -U/--update is used.')
    arg_parser.add_argument('-d', '--data_file', required=True, help='CSV file containing the chart data and format specifications.')
    arg_parser.add_argument('-i', '--input_file', help='Existing pptx file to add the chart(s) to or to update if used with -U/--update.')
    arg_parser.add_argument('-s', '--slide', type=int, default=1, help='Index of the slide to modify.')
    arg_parser.add_argument('-k', '--shape_id', help='Selection pane name for the shape that contains the chart to be updated (use with -U/--update).')
    arg_parser.add_argument('-U', '--update', action='store_true', help='Update an existing chart that has a selection pane name equal to -k/--shape_id or the value of the chart.id column (only if the facet.col is present).')
    arg_parser.add_argument('--update-format', action='store_true', help='Update chart formats (use with --U/--update).')
    arg_parser.add_argument('--ignore-missing-charts', action='store_true', help='Continue without raising an error if the chart cannot be found (use with -U/--update).')
    args = arg_parser.parse_args()
    slide = (args.slide - 1) if args.slide is not None else None

    if not args.update:
        if args.output_file is None:
            raise ValueError('Argument -o/--output_file is required')
        add_chart(
            output_file=args.output_file,
            data_file=args.data_file,
            slide_idx=slide,
            input_file=args.input_file
        )
    else:
        if args.slide is None:
            raise ValueError('Argument -s/--slide is required when using -U/--update')
        if args.input_file is None:
            raise ValueError('Argument -i/--input_file is required when using -U/--update')
        update_chart(
            output_file=args.output_file,
            data_file=args.data_file,
            slide_idx=slide,
            input_file=args.input_file,
            shape_id=args.shape_id,
            should_update_format=args.update_format,
            ignore_missing_charts=args.ignore_missing_charts
        )


if __name__ == '__main__':
    main()
